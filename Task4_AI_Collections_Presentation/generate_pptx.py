"""
Task 4: AI-Powered Collections System — Executive Briefing PowerPoint
Geldium | Tata iQ Analytics Team | March 2026
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color palette ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
GRAY_TEXT = RGBColor(0x7F, 0x8C, 0x8D)
GREEN = RGBColor(0x27, 0xAE, 0x60)
ORANGE = RGBColor(0xE6, 0x7E, 0x22)
RED = RGBColor(0xC0, 0x39, 0x2B)
LIGHT_BLUE_BG = RGBColor(0xD6, 0xEA, 0xF8)
LIGHT_GREEN_BG = RGBColor(0xD5, 0xF5, 0xE3)
LIGHT_ORANGE_BG = RGBColor(0xFD, 0xED, 0xEC)
LIGHT_PURPLE_BG = RGBColor(0xEB, 0xDE, 0xF0)

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape_box(slide, left, top, width, height, fill_color, text="", font_size=12, bold=False, font_color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_right = Inches(0.2)
        tf.margin_top = Inches(0.15)
        tf.margin_bottom = Inches(0.1)
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.color.rgb = font_color
        p.alignment = alignment
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return tf

def add_bullet_frame(slide, left, top, width, height, bullets, font_size=13, color=DARK_TEXT, bold_prefix=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.level = 0
        if isinstance(bullet, tuple) and bold_prefix:
            run = p.add_run()
            run.text = bullet[0]
            run.font.bold = True
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run = p.add_run()
            run.text = bullet[1]
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
        else:
            txt = bullet if isinstance(bullet, str) else bullet[0] + bullet[1]
            run = p.add_run()
            run.text = txt
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return tf

def add_arrow(slide, left, top, width, height):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_BLUE
    shape.line.fill.background()
    return shape

# ══════════════════════════════════════════════════════════════
# SLIDE 0: TITLE SLIDE
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(1.2),
             "AI-Powered Collections System", font_size=40, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10.3), Inches(0.8),
             "A Framework for Autonomous, Responsible Debt Management", font_size=22, color=RGBColor(0xAE, 0xD6, 0xF1), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1.5), Inches(4.5), Inches(10.3), Inches(0.6),
             "Executive Briefing for Geldium | Tata iQ Analytics Team | March 2026", font_size=14, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# Key stats bar
for i, (label, value) in enumerate([("Customers Analyzed", "500"), ("Delinquency Rate", "16%"), ("Top Predictor", "Payment Behavior"), ("Model AUC-ROC", "0.93")]):
    x = Inches(1.8) + Inches(2.5) * i
    box = add_shape_box(slide, x, Inches(5.5), Inches(2.2), Inches(1.0), RGBColor(0x23, 0x4E, 0x73), alignment=PP_ALIGN.CENTER)
    tf = box.text_frame
    tf.clear()
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = label
    run2.font.size = Pt(10)
    run2.font.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
    p2.alignment = PP_ALIGN.CENTER

# ══════════════════════════════════════════════════════════════
# SLIDE 1: HOW THE SYSTEM WORKS (Overview)
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

# Title bar
add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), DARK_BLUE,
              "  System Overview: How It Works", font_size=26, bold=True, font_color=WHITE)

add_text_box(slide, Inches(0.6), Inches(1.1), Inches(12), Inches(0.5),
             "The AI-powered collections system operates as a continuous 4-stage loop, transforming raw customer data into personalized, adaptive outreach.",
             font_size=13, color=GRAY_TEXT)

# 4-part flow diagram
stages = [
    ("1. DATA PIPELINE", "Ingest & Enrich", [
        "Real-time customer data feeds",
        "Risk scores from predictive model",
        "Payment history & credit bureau",
        "Behavioral signals & trends"
    ], LIGHT_BLUE_BG, ACCENT_BLUE),
    ("2. DECISION ENGINE", "Analyze & Decide", [
        "Gradient Boosting risk scoring",
        "Segment into Low / Medium / High",
        "Match risk level to intervention",
        "Apply business rules & constraints"
    ], LIGHT_GREEN_BG, GREEN),
    ("3. ACTION LAYER", "Execute & Personalize", [
        "Trigger SMS/email reminders",
        "Offer payment deferrals",
        "Adjust outreach timing & tone",
        "Escalate high-risk to human agents"
    ], RGBColor(0xFE, 0xF9, 0xE7), ORANGE),
    ("4. LEARNING LOOP", "Monitor & Improve", [
        "Track repayment outcomes",
        "Measure engagement rates",
        "Retrain model quarterly",
        "Update strategies from feedback"
    ], LIGHT_PURPLE_BG, RGBColor(0x8E, 0x44, 0xAD)),
]

box_w = Inches(2.7)
box_h = Inches(4.2)
start_x = Inches(0.5)
gap = Inches(0.55)

for i, (title, subtitle, bullets, bg_color, accent) in enumerate(stages):
    x = start_x + (box_w + gap) * i
    y = Inches(1.8)

    # Main box
    box = add_shape_box(slide, x, y, box_w, box_h, bg_color)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)

    # Stage title
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(14)
    run.font.bold = True
    run.font.color.rgb = accent
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.size = Pt(11)
    run2.font.italic = True
    run2.font.color.rgb = GRAY_TEXT
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(10)

    # Bullets
    for b in bullets:
        pb = tf.add_paragraph()
        run_b = pb.add_run()
        run_b.text = "  " + b
        run_b.font.size = Pt(11)
        run_b.font.color.rgb = DARK_TEXT
        pb.space_before = Pt(3)
        pb.space_after = Pt(2)

    # Arrow between boxes
    if i < 3:
        arrow_x = x + box_w + Inches(0.05)
        add_arrow(slide, arrow_x, Inches(3.6), Inches(0.45), Inches(0.4))

# Circular arrow annotation
add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.6),
             "The Learning Loop feeds insights back into the Data Pipeline, creating a self-improving cycle.",
             font_size=12, bold=True, color=ACCENT_BLUE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 2: SYSTEM WORKFLOW DETAIL
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), DARK_BLUE,
              "  System Workflow: From Risk Score to Action", font_size=26, bold=True, font_color=WHITE)

# Left: Risk tiers
add_text_box(slide, Inches(0.5), Inches(1.2), Inches(4), Inches(0.4),
             "Risk-Based Decision Logic", font_size=16, bold=True, color=DARK_BLUE)

tiers = [
    ("LOW RISK (Score 0 - 0.3)", "430 customers (86%)", [
        "Standard monthly statement reminders",
        "Self-service payment portal nudges",
        "No direct Collections contact needed"
    ], LIGHT_GREEN_BG, GREEN),
    ("MEDIUM RISK (Score 0.3 - 0.6)", "55 customers (11%)", [
        "Personalized SMS/email reminders",
        "Proactive payment plan offers",
        "Automated hardship screening",
    ], RGBColor(0xFE, 0xF9, 0xE7), ORANGE),
    ("HIGH RISK (Score 0.6 - 1.0)", "15 customers (3%)", [
        "Immediate human agent assignment",
        "Tailored repayment restructuring",
        "Escalation to senior Collections team"
    ], LIGHT_ORANGE_BG, RED),
]

for i, (title, count, actions, bg, accent) in enumerate(tiers):
    y = Inches(1.8) + Inches(1.7) * i
    box = add_shape_box(slide, Inches(0.5), y, Inches(5.8), Inches(1.5), bg)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = accent

    run2 = p.add_run()
    run2.text = f"  |  {count}"
    run2.font.size = Pt(11)
    run2.font.color.rgb = GRAY_TEXT

    for a in actions:
        pa = tf.add_paragraph()
        ra = pa.add_run()
        ra.text = "    " + a
        ra.font.size = Pt(11)
        ra.font.color.rgb = DARK_TEXT
        pa.space_before = Pt(2)

# Right: Key model inputs
add_text_box(slide, Inches(7), Inches(1.2), Inches(5.5), Inches(0.4),
             "Key Model Inputs (Top 5 Predictors)", font_size=16, bold=True, color=DARK_BLUE)

predictors = [
    ("Avg Payment Severity", "12.3%", "6-month payment behavior pattern"),
    ("Month 5 Status", "9.1%", "Recent payment — strong near-term signal"),
    ("Month 2 Status", "8.9%", "Early behavior — establishes baseline"),
    ("Month 6 Status", "7.5%", "Most recent — current risk indicator"),
    ("Credit Utilization", "6.3%", "Financial strain signal"),
]

# Table
table_shape = slide.shapes.add_table(len(predictors) + 1, 3, Inches(7), Inches(1.8), Inches(5.8), Inches(2.5))
table = table_shape.table
table.columns[0].width = Inches(2.2)
table.columns[1].width = Inches(1.2)
table.columns[2].width = Inches(2.4)

headers = ["Feature", "Importance", "Why It Matters"]
for ci, h in enumerate(headers):
    cell = table.cell(0, ci)
    cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
    cell.fill.solid()
    cell.fill.fore_color.rgb = DARK_BLUE

for ri, (feat, imp, why) in enumerate(predictors):
    for ci, val in enumerate([feat, imp, why]):
        cell = table.cell(ri + 1, ci)
        cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10.5)
            p.font.color.rgb = DARK_TEXT
        if ri % 2 == 0:
            cell.fill.solid()
            cell.fill.fore_color.rgb = LIGHT_BG

# Real-time adaptation box
add_shape_box(slide, Inches(7), Inches(4.8), Inches(5.8), Inches(1.8), LIGHT_BLUE_BG)
tf = add_text_box(slide, Inches(7.2), Inches(4.9), Inches(5.4), Inches(1.6),
                  "Real-Time Adaptation Example", font_size=13, bold=True, color=ACCENT_BLUE)
p2 = tf.add_paragraph()
run = p2.add_run()
run.text = ('A customer currently scored "Low Risk" misses two consecutive payments. '
            'The system automatically re-scores them, upgrades their risk tier to "Medium," '
            'and triggers a personalized SMS reminder with a payment plan offer — all within '
            '24 hours, without manual intervention.')
run.font.size = Pt(11)
run.font.color.rgb = DARK_TEXT
p2.space_before = Pt(6)

# ══════════════════════════════════════════════════════════════
# SLIDE 3: ROLE OF AGENTIC AI
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), DARK_BLUE,
              "  Role of Agentic AI: Autonomous vs. Human Oversight", font_size=26, bold=True, font_color=WHITE)

add_text_box(slide, Inches(0.5), Inches(1.1), Inches(12), Inches(0.5),
             "Agentic AI makes autonomous decisions where speed and scale matter, while human oversight is preserved for high-stakes and sensitive actions.",
             font_size=13, color=GRAY_TEXT)

# Two-column table
table_shape = slide.shapes.add_table(7, 2, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.5))
table = table_shape.table
table.columns[0].width = Inches(6.15)
table.columns[1].width = Inches(6.15)

# Headers
for ci, (h, color) in enumerate([("Autonomous AI Actions", GREEN), ("Human-in-the-Loop Oversight", ORANGE)]):
    cell = table.cell(0, ci)
    cell.text = h
    cell.fill.solid()
    cell.fill.fore_color.rgb = color
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

auto_items = [
    "Real-time risk scoring and tier assignment for all 500+ customers",
    "Triggering personalized SMS/email reminders based on risk level and behavior",
    "Dynamically adjusting outreach timing and frequency based on engagement data",
    "Re-scoring customers when new payment data arrives (within 24 hours)",
    "Routing medium-risk cases to automated hardship screening workflows",
    "Monitoring repayment outcomes and updating model confidence scores",
]
human_items = [
    "Approving or denying hardship assistance and payment restructuring requests",
    "Reviewing all high-risk escalations before legal or final collections action",
    "Overriding AI recommendations when individual circumstances warrant exceptions",
    "Quarterly model performance review and bias audit sign-off",
    "Setting and adjusting business rules, thresholds, and outreach tone guidelines",
    "Final approval on any action that could negatively impact a customer's credit record",
]

for ri, (auto, human) in enumerate(zip(auto_items, human_items)):
    bg = LIGHT_BG if ri % 2 == 0 else WHITE
    for ci, text in enumerate([auto, human]):
        cell = table.cell(ri + 1, ci)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11.5)
            p.font.color.rgb = DARK_TEXT

# Bottom callout
add_shape_box(slide, Inches(1.5), Inches(6.5), Inches(10.3), Inches(0.7), LIGHT_BLUE_BG,
              "Principle: AI handles speed and scale. Humans handle judgment and accountability. No customer-harming action is taken without human review.",
              font_size=12, bold=True, font_color=DARK_BLUE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 4: RESPONSIBLE AI GUARDRAILS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), DARK_BLUE,
              "  Responsible AI Guardrails", font_size=26, bold=True, font_color=WHITE)

guardrails = [
    ("Fairness & Bias Prevention", LIGHT_GREEN_BG, GREEN, [
        "Disparate impact audits across employment, location, and age groups every quarter",
        "Exclude proxy variables (e.g., location) if they encode demographic bias without improving accuracy",
        "Group-specific threshold calibration to ensure predicted rates align with actual rates",
        "Diverse, representative training data with SMOTE to prevent class imbalance bias"
    ]),
    ("Explainability & Transparency", LIGHT_BLUE_BG, ACCENT_BLUE, [
        "SHAP values provide per-customer explanations: 'Flagged due to 3 missed payments + 92% credit utilization'",
        "Companion Decision Tree model generates plain-language if/then rules for Collections staff",
        "Full audit trail of every AI decision — what was decided, why, and what data was used",
        "Customer-facing explanations for any adverse action, with clear appeal process"
    ]),
    ("Regulatory Compliance", RGBColor(0xFE, 0xF9, 0xE7), ORANGE, [
        "ECOA: Model validated to not discriminate by protected characteristics",
        "GDPR: Right to explanation honored — customers can request reasons for any automated decision",
        "FCA: Proportionate collections — system prevents overly aggressive outreach",
        "FCRA: Only current, accurate data used; outdated records excluded from scoring"
    ]),
    ("Continuous Monitoring", LIGHT_PURPLE_BG, RGBColor(0x8E, 0x44, 0xAD), [
        "Quarterly model retraining with fresh customer data to prevent drift",
        "Monthly fairness metric dashboards reviewed by Compliance team",
        "Automated alerts if Disparate Impact Ratio falls below 0.80 threshold",
        "Annual independent audit by external ethics review board"
    ]),
]

box_w = Inches(6.0)
box_h = Inches(2.6)
positions = [
    (Inches(0.4), Inches(1.2)),
    (Inches(6.9), Inches(1.2)),
    (Inches(0.4), Inches(4.0)),
    (Inches(6.9), Inches(4.0)),
]

for i, ((title, bg, accent, bullets), (x, y)) in enumerate(zip(guardrails, positions)):
    box = add_shape_box(slide, x, y, box_w, box_h, bg)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.12)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = accent
    p.space_after = Pt(6)

    for b in bullets:
        pb = tf.add_paragraph()
        rb = pb.add_run()
        rb.text = "  " + b
        rb.font.size = Pt(10.5)
        rb.font.color.rgb = DARK_TEXT
        pb.space_before = Pt(3)
        pb.space_after = Pt(1)

# ══════════════════════════════════════════════════════════════
# SLIDE 5: EXPECTED BUSINESS IMPACT
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)

add_shape_box(slide, Inches(0), Inches(0), Inches(13.333), Inches(0.9), DARK_BLUE,
              "  Expected Business Impact", font_size=26, bold=True, font_color=WHITE)

# Quantitative — left side
add_text_box(slide, Inches(0.5), Inches(1.1), Inches(6), Inches(0.4),
             "Quantitative Outcomes", font_size=18, bold=True, color=GREEN)

quant_metrics = [
    ("15%", "Reduction in delinquency rate\namong high-risk segments", GREEN),
    ("40%", "Reduction in manual outreach\neffort through automation", ACCENT_BLUE),
    ("25%", "Improvement in early intervention\nsuccess rate", ORANGE),
    ("20%", "Reduction in cost-per-collection\nthrough targeted outreach", RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (number, desc, color) in enumerate(quant_metrics):
    x = Inches(0.5) + Inches(3.1) * (i % 2)
    y = Inches(1.7) + Inches(1.5) * (i // 2)

    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(1.0), Inches(1.0))
    circle.fill.solid()
    circle.fill.fore_color.rgb = color
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = number
    run.font.size = Pt(22)
    run.font.bold = True
    run.font.color.rgb = WHITE
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Description
    add_text_box(slide, x + Inches(1.15), y + Inches(0.1), Inches(1.85), Inches(0.9),
                 desc, font_size=11, color=DARK_TEXT)

# Qualitative — right side
add_text_box(slide, Inches(6.9), Inches(1.1), Inches(6), Inches(0.4),
             "Qualitative Outcomes", font_size=18, bold=True, color=ACCENT_BLUE)

qual_items = [
    ("Better Customer Experience",
     "Personalized, empathetic outreach replaces generic collection calls. "
     "Customers receive support tailored to their financial situation, improving satisfaction and retention."),
    ("Increased Fairness & Trust",
     "Transparent AI decisions with SHAP explanations and human oversight build customer "
     "and regulator confidence. No demographic group is systematically disadvantaged."),
    ("Operational Scalability",
     "The system handles growing customer volumes without proportional staff increases. "
     "Collections team focuses on complex cases while AI manages routine outreach."),
    ("Regulatory Confidence",
     "Built-in compliance with ECOA, GDPR, FCA, and FCRA standards ensures Geldium "
     "is audit-ready and protected against regulatory penalties."),
]

for i, (title, desc) in enumerate(qual_items):
    y = Inches(1.7) + Inches(1.3) * i
    box = add_shape_box(slide, Inches(6.9), y, Inches(5.9), Inches(1.15), LIGHT_BG)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(13)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(10.5)
    run2.font.color.rgb = DARK_TEXT
    p2.space_before = Pt(4)

# Bottom summary bar
add_shape_box(slide, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), DARK_BLUE,
              "  Bottom Line: AI-powered collections delivers measurable financial savings, fairer outcomes, and a scalable system Geldium can grow with.",
              font_size=13, bold=True, font_color=WHITE, alignment=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════
# SLIDE 6: SUMMARY & NEXT STEPS
# ══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)

add_text_box(slide, Inches(1), Inches(0.6), Inches(11.3), Inches(0.6),
             "Summary & Next Steps", font_size=30, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

# Summary boxes
summaries = [
    ("System Design", "4-stage continuous loop:\nData Pipeline > Decision Engine >\nAction Layer > Learning Loop"),
    ("Agentic AI", "Autonomous for speed & scale.\nHuman oversight for judgment\n& high-stakes decisions."),
    ("Guardrails", "Fairness audits, SHAP explainability,\nregulatory compliance (ECOA/GDPR/\nFCA), continuous monitoring."),
    ("Business Impact", "15% delinquency reduction,\n40% less manual effort,\nfairer & scalable collections."),
]

for i, (title, desc) in enumerate(summaries):
    x = Inches(0.5) + Inches(3.2) * i
    box = add_shape_box(slide, x, Inches(1.6), Inches(2.9), Inches(2.2), RGBColor(0x23, 0x4E, 0x73))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_top = Inches(0.15)

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(15)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0xAE, 0xD6, 0xF1)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(10)

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = desc
    run2.font.size = Pt(11.5)
    run2.font.color.rgb = WHITE
    p2.alignment = PP_ALIGN.CENTER

# Next steps
add_text_box(slide, Inches(1), Inches(4.2), Inches(11.3), Inches(0.5),
             "Recommended Next Steps", font_size=20, bold=True, color=RGBColor(0xAE, 0xD6, 0xF1), alignment=PP_ALIGN.CENTER)

steps = [
    "1. Pilot Phase (90 days): Deploy AI risk scoring for Medium-Risk segment with automated SMS reminders",
    "2. Validation (Month 4-6): Measure delinquency reduction, fairness metrics, and customer satisfaction",
    "3. Scale (Month 7-12): Expand to all risk tiers, integrate with CRM, and enable full adaptive learning loop",
    "4. Ongoing: Quarterly model retraining, monthly fairness audits, annual independent ethics review",
]

for i, step in enumerate(steps):
    y = Inches(4.8) + Inches(0.45) * i
    add_text_box(slide, Inches(1.5), y, Inches(10.3), Inches(0.4),
                 step, font_size=12, color=WHITE)

# AI usage note
add_text_box(slide, Inches(1), Inches(6.8), Inches(11.3), Inches(0.4),
             "AI & GenAI Disclosure: This presentation was developed with assistance from Claude Code (GenAI) for analysis, framework design, and slide generation.",
             font_size=9, color=GRAY_TEXT, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = 'C:/Users/LENOVO X1 YOGA/Desktop/tata/Geldium_AI_Collections_System.pptx'
prs.save(output_path)
print(f'Presentation saved to: {output_path}')
print(f'Total slides: {len(prs.slides)}')
